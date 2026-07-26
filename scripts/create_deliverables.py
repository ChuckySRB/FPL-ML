"""Build the PDF report and five-minute PowerPoint presentation."""

from __future__ import annotations

import re
from pathlib import Path
from xml.sax.saxutils import escape

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    Image,
    ListFlowable,
    ListItem,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


ROOT = Path(__file__).resolve().parents[1]
REPORT_MD = ROOT / "reports" / "final_report.md"
REPORT_PDF = ROOT / "reports" / "final_report.pdf"
PRESENTATION_PPTX = ROOT / "presentation" / "final_presentation.pptx"
FIGURES = ROOT / "outputs" / "figures"

NAVY = RGBColor(12, 30, 56)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
LIGHT = RGBColor(241, 245, 249)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(30, 41, 59)


def register_pdf_fonts() -> tuple[str, str]:
    """Register a Windows font with Cyrillic support."""
    choices = [
        (Path("C:/Windows/Fonts/arial.ttf"), Path("C:/Windows/Fonts/arialbd.ttf")),
        (
            Path("C:/Windows/Fonts/DejaVuSans.ttf"),
            Path("C:/Windows/Fonts/DejaVuSans-Bold.ttf"),
        ),
    ]
    for normal, bold in choices:
        if normal.exists() and bold.exists():
            pdfmetrics.registerFont(TTFont("ProjectSans", str(normal)))
            pdfmetrics.registerFont(TTFont("ProjectSans-Bold", str(bold)))
            return "ProjectSans", "ProjectSans-Bold"
    return "Helvetica", "Helvetica-Bold"


def inline_markup(text: str) -> str:
    """Convert the small Markdown subset used by the report."""
    text = escape(text)
    text = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", text)
    text = re.sub(r"`(.+?)`", r"<font name='Courier'>\1</font>", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1 (\2)", text)
    return text


def scaled_report_image(path: Path, max_width: float = 16.5 * cm) -> Image:
    """Create a proportionally scaled report image."""
    with PILImage.open(path) as image:
        width, height = image.size
    ratio = min(max_width / width, 10.5 * cm / height)
    return Image(str(path), width=width * ratio, height=height * ratio)


def report_footer(canvas, doc) -> None:
    """Draw page number and report label."""
    canvas.saveState()
    canvas.setFont("ProjectSans", 8)
    canvas.setFillColor(colors.HexColor("#64748b"))
    canvas.drawString(2 * cm, 1.15 * cm, "OpenFPL — делимична репродукција")
    canvas.drawRightString(19 * cm, 1.15 * cm, f"Страна {doc.page}")
    canvas.restoreState()


def build_report() -> None:
    """Render the Markdown report to a submission-ready PDF."""
    font, bold_font = register_pdf_fonts()
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "Body",
        parent=styles["BodyText"],
        fontName=font,
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#1e293b"),
        spaceAfter=7,
    )
    h1 = ParagraphStyle(
        "H1",
        parent=body,
        fontName=bold_font,
        fontSize=22,
        leading=27,
        alignment=TA_CENTER,
        textColor=colors.HexColor("#0c1e38"),
        spaceAfter=20,
    )
    h2 = ParagraphStyle(
        "H2",
        parent=body,
        fontName=bold_font,
        fontSize=15,
        leading=19,
        textColor=colors.HexColor("#2563eb"),
        spaceBefore=10,
        spaceAfter=7,
        keepWithNext=True,
    )
    bullet_style = ParagraphStyle("Bullet", parent=body, leftIndent=8, spaceAfter=2)
    table_style = ParagraphStyle(
        "TableCell", parent=body, fontSize=8.3, leading=10, spaceAfter=0
    )

    document = SimpleDocTemplate(
        str(REPORT_PDF),
        pagesize=A4,
        rightMargin=1.8 * cm,
        leftMargin=1.8 * cm,
        topMargin=1.7 * cm,
        bottomMargin=1.8 * cm,
        title="Предвиђање учинка играча у Fantasy Premier League-у",
        author="Машинско учење — пројекат",
    )
    lines = REPORT_MD.read_text(encoding="utf-8").splitlines()
    story = []
    paragraph_lines: list[str] = []

    def flush_paragraph() -> None:
        if paragraph_lines:
            story.append(Paragraph(inline_markup(" ".join(paragraph_lines)), body))
            paragraph_lines.clear()

    index = 0
    while index < len(lines):
        line = lines[index].strip()
        if not line:
            flush_paragraph()
            index += 1
            continue
        if line.startswith("# "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[2:]), h1))
        elif line.startswith("## "):
            flush_paragraph()
            story.append(Paragraph(inline_markup(line[3:]), h2))
        elif line.startswith("!["):
            flush_paragraph()
            match = re.match(r"!\[(.*?)\]\((.*?)\)", line)
            if match:
                image_path = (REPORT_MD.parent / match.group(2)).resolve()
                story.append(scaled_report_image(image_path))
                story.append(
                    Paragraph(
                        f"<i>{escape(match.group(1))}</i>",
                        ParagraphStyle(
                            "Caption", parent=body, fontSize=8, alignment=TA_CENTER
                        ),
                    )
                )
                story.append(Spacer(1, 5))
        elif line.startswith("|") and index + 1 < len(lines):
            flush_paragraph()
            table_lines = [line]
            index += 1
            while index < len(lines) and lines[index].strip().startswith("|"):
                table_lines.append(lines[index].strip())
                index += 1
            rows = []
            for row_index, table_line in enumerate(table_lines):
                cells = [cell.strip() for cell in table_line.strip("|").split("|")]
                if row_index == 1 and all(set(cell) <= set("-: ") for cell in cells):
                    continue
                rows.append([Paragraph(inline_markup(cell), table_style) for cell in cells])
            table = Table(rows, repeatRows=1, hAlign="LEFT")
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0c1e38")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("FONTNAME", (0, 0), (-1, 0), bold_font),
                        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#cbd5e1")),
                        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [
                            colors.white,
                            colors.HexColor("#f8fafc"),
                        ]),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("LEFTPADDING", (0, 0), (-1, -1), 5),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                        ("TOPPADDING", (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ]
                )
            )
            story.extend([table, Spacer(1, 7)])
            continue
        elif re.match(r"^\d+\.\s", line):
            flush_paragraph()
            items = []
            while index < len(lines) and re.match(r"^\d+\.\s", lines[index].strip()):
                item = re.sub(r"^\d+\.\s*", "", lines[index].strip())
                items.append(ListItem(Paragraph(inline_markup(item), bullet_style)))
                index += 1
            story.append(ListFlowable(items, bulletType="1", leftIndent=18))
            continue
        elif line.startswith("- "):
            flush_paragraph()
            items = []
            while index < len(lines) and lines[index].strip().startswith("- "):
                items.append(
                    ListItem(
                        Paragraph(inline_markup(lines[index].strip()[2:]), bullet_style)
                    )
                )
                index += 1
            story.append(ListFlowable(items, bulletType="bullet", leftIndent=18))
            continue
        else:
            paragraph_lines.append(line)
        index += 1
    flush_paragraph()
    document.build(story, onFirstPage=report_footer, onLaterPages=report_footer)


def add_textbox(slide, left, top, width, height, text, size=20, color=DARK,
                bold=False, align=PP_ALIGN.LEFT):
    """Add a consistently styled PowerPoint text box."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def base_slide(prs: Presentation, title: str):
    """Create a blank slide with the shared header/footer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    add_textbox(
        slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.65),
        title, size=27, color=NAVY, bold=True
    )
    add_textbox(
        slide, Inches(10.5), Inches(7.05), Inches(2.25), Inches(0.22),
        "OpenFPL reproduction", size=8, color=RGBColor(100, 116, 139),
        align=PP_ALIGN.RIGHT
    )
    return slide


def add_picture_contain(slide, path: Path, left, top, width, height) -> None:
    """Place an image without distortion inside the requested rectangle."""
    with PILImage.open(path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio > box_ratio:
        draw_width = width
        draw_height = width / image_ratio
        draw_left = left
        draw_top = top + (height - draw_height) / 2
    else:
        draw_height = height
        draw_width = height * image_ratio
        draw_top = top
        draw_left = left + (width - draw_width) / 2
    slide.shapes.add_picture(str(path), draw_left, draw_top, draw_width, draw_height)


def add_card(slide, left, top, width, height, heading, value, color=BLUE):
    """Add a metric card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(203, 213, 225)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.08),
                width - Inches(0.32), Inches(0.3), heading, 11,
                RGBColor(71, 85, 105), True, PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.42),
                width - Inches(0.24), height - Inches(0.5), value, 23,
                color, True, PP_ALIGN.CENTER)


def add_notes(slide, text: str) -> None:
    """Attach the talk track to the PowerPoint notes pane."""
    frame = slide.notes_slide.notes_text_frame
    if frame is not None:
        frame.text = text


def build_presentation() -> None:
    """Create a visual six-slide presentation with speaker notes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.0), Inches(0.14), Inches(4.7)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()
    add_textbox(
        slide, Inches(1.15), Inches(1.35), Inches(10.8), Inches(1.6),
        "Предвиђање FPL поена", 38, WHITE, True
    )
    add_textbox(
        slide, Inches(1.18), Inches(3.0), Inches(10.6), Inches(0.8),
        "Leakage-safe ML pipeline и делимична репродукција OpenFPL рада",
        21, RGBColor(191, 219, 254)
    )
    add_textbox(
        slide, Inches(1.18), Inches(5.65), Inches(10.4), Inches(0.45),
        "Машинско учење · 2026", 14, RGBColor(148, 163, 184)
    )
    add_notes(slide, "Циљ је да за сваког играча пре утакмице проценим број "
              "FPL поена. Главни циљ био је временски исправан тест без "
              "цурења података и поређење са OpenFPL радом.")

    slide = base_slide(prs, "Подаци: много нула, ретки велики учинци")
    add_picture_contain(
        slide, FIGURES / "eda_target_distribution_fixed.png",
        Inches(0.55), Inches(1.15), Inches(8.1), Inches(5.55)
    )
    add_card(slide, Inches(9.0), Inches(1.45), Inches(3.5), Inches(1.2),
             "TRAIN", "94.164 реда")
    add_card(slide, Inches(9.0), Inches(2.95), Inches(3.5), Inches(1.2),
             "TEST 2024/25", "24.394 реда", GREEN)
    add_card(slide, Inches(9.0), Inches(4.45), Inches(3.5), Inches(1.2),
             "PROSPECTIVE", "GW32–38", RGBColor(124, 58, 237))
    add_textbox(
        slide, Inches(9.0), Inches(5.9), Inches(3.5), Inches(0.6),
        "Јединица: играч–утакмица", 14, DARK, True, PP_ALIGN.CENTER
    )
    add_notes(slide, "Тренирао сам на сезонама 2020/21–2023/24, а 2024/25 "
              "оставио само за тест. Расподела има много нула и дугачак реп, "
              "па гледам RMSE и категорије, не само MAE.")

    slide = base_slide(prs, "Методологија: прошлост не сме видети будућност")
    labels = [
        ("5 сезона\nFPL података", BLUE),
        ("shift(1)\nrolling features", GREEN),
        ("9 модела\n+ baseline-и", RGBColor(124, 58, 237)),
        ("2024/25\nGW32–38", RGBColor(220, 38, 38)),
    ]
    lefts = [0.7, 3.85, 7.0, 10.15]
    for index, ((label, color), left) in enumerate(zip(labels, lefts)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.1),
            Inches(2.45), Inches(1.75)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        add_textbox(
            slide, Inches(left + 0.15), Inches(2.25), Inches(2.15), Inches(1.4),
            label, 19, WHITE, True, PP_ALIGN.CENTER
        )
        if index < 3:
            add_textbox(
                slide, Inches(left + 2.47), Inches(2.55), Inches(0.65),
                Inches(0.65), "→", 27, NAVY, True, PP_ALIGN.CENTER
            )
    add_textbox(
        slide, Inches(1.0), Inches(4.65), Inches(11.2), Inches(1.0),
        "DGW-safe · fixture-specific FDR · median imputation learned on train · stable-club CV",
        17, DARK, True, PP_ALIGN.CENTER
    )
    add_notes(slide, "Сва обележја користе shift један, одвојена су по "
              "сезони и играчу, а double gameweek не види исход другог меча. "
              "Поређени су једноставни baseline-и, Tier-0, Linear, RF и XGBoost.")

    slide = base_slide(prs, "Резултат: три најбоља модела су врло близу")
    add_picture_contain(
        slide, FIGURES / "model_comparison_fixed.png",
        Inches(0.45), Inches(1.0), Inches(9.0), Inches(5.9)
    )
    add_card(slide, Inches(9.7), Inches(1.45), Inches(2.9), Inches(1.2),
             "XGBOOST RMSE", "1,842", GREEN)
    add_card(slide, Inches(9.7), Inches(2.95), Inches(2.9), Inches(1.2),
             "RF RMSE", "1,846")
    add_card(slide, Inches(9.7), Inches(4.45), Inches(2.9), Inches(1.2),
             "LINEAR RMSE", "1,856", RGBColor(124, 58, 237))
    add_textbox(
        slide, Inches(9.7), Inches(5.85), Inches(2.9), Inches(0.65),
        "Сложеније ≠ боље", 16, RGBColor(220, 38, 38), True, PP_ALIGN.CENTER
    )
    add_notes(slide, "На GW32–38 најбољи је необтежени XGBoost. RF и "
              "линеарна регресија су веома близу. Тежински и position ансамбл "
              "су лошији, што је важан негативан резултат.")

    slide = base_slide(prs, "OpenFPL поређење и главни failure mode")
    add_picture_contain(
        slide, FIGURES / "category_rmse_fixed.png",
        Inches(0.45), Inches(1.0), Inches(7.65), Inches(5.9)
    )
    add_textbox(
        slide, Inches(8.35), Inches(1.25), Inches(4.35), Inches(0.55),
        "Највећи промашаји", 20, NAVY, True, PP_ALIGN.CENTER
    )
    errors = [
        "Mbeumo: 18 → 4,73",
        "Colwill: 15 → 1,80",
        "Ramsdale: 16 → 3,08",
    ]
    for index, text in enumerate(errors):
        add_card(
            slide, Inches(8.55), Inches(2.0 + index * 1.15),
            Inches(3.95), Inches(0.88), f"#{index + 1}", text,
            RGBColor(220, 38, 38)
        )
    add_textbox(
        slide, Inches(8.45), Inches(5.65), Inches(4.05), Inches(0.85),
        "Hauler-и су систематски потцењени", 16,
        RGBColor(220, 38, 38), True, PP_ALIGN.CENTER
    )
    add_notes(slide, "У односу на OpenFPL рад бољи сам за Zeros и Tickers, "
              "а слабији за Blanks и Haulers. Поређење је делимично јер немам "
              "сва Understat и availability обележја. Највеће грешке су ретки "
              "голови, пенали и clean sheet учинци.")

    slide = base_slide(prs, "Закључак и следећа надоградња")
    conclusions = [
        ("ПОУЗДАН PIPELINE", "Временски split\nбез leakage-а", BLUE),
        ("НАЈБОЉИ МОДЕЛ", "XGBoost\nRMSE 1,842", GREEN),
        ("ОТВОРЕН ПРОБЛЕМ", "Неизвесност\nи hauler-и", RGBColor(220, 38, 38)),
    ]
    for index, (heading, value, color) in enumerate(conclusions):
        add_card(
            slide, Inches(0.75 + index * 4.2), Inches(1.45),
            Inches(3.75), Inches(1.7), heading, value, color
        )
    add_textbox(
        slide, Inches(0.85), Inches(4.1), Inches(11.7), Inches(0.55),
        "Следеће: fresh availability → prediction intervals → оптимизатор тима",
        20, NAVY, True, PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(1.2), Inches(5.2), Inches(10.9), Inches(0.7),
        "Најважнија лекција: јак baseline и исправна евалуација "
        "вреде више од саме сложености.",
        18, RGBColor(71, 85, 105), False, PP_ALIGN.CENTER
    )
    add_notes(slide, "Најважнији резултат је поуздан временски pipeline и "
              "јасна baseline лествица. XGBoost је најбољи, али добитак је "
              "мали. Следеће су prediction интервали, свежа availability "
              "обележја и практичан оптимизатор.")

    prs.save(PRESENTATION_PPTX)


def main() -> None:
    """Build both final deliverables."""
    REPORT_PDF.parent.mkdir(parents=True, exist_ok=True)
    PRESENTATION_PPTX.parent.mkdir(parents=True, exist_ok=True)
    build_report()
    build_presentation()
    print("Created final report PDF and PowerPoint presentation")


if __name__ == "__main__":
    main()
